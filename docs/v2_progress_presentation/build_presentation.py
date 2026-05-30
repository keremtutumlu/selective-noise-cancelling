"""
Build the v2.x iteration progress presentation as a .pptx file.

Run:
    python docs/v2_progress_presentation/build_presentation.py

Output:
    docs/v2_progress_presentation/presentation.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).with_name("presentation.pptx")

# --- Theme ---------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x7A, 0x7A, 0x7A)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x1F, 0x8A, 0x4C)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HDR = NAVY
TABLE_ALT = RGBColor(0xF2, 0xF5, 0xFA)

BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

# Slide is 16:9 (13.333 x 7.5 inches)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# --- helpers -------------------------------------------------------------
def add_title_bar(slide, title_text):
    """Navy bar at the top, white title text."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = BODY_FONT
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_footer(slide, page_num, total):
    """Slide number bottom-right."""
    tb = slide.shapes.add_textbox(prs.slide_width - Inches(2.0),
                                   prs.slide_height - Inches(0.4),
                                   Inches(1.8), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num} / {total}"
    run.font.name = BODY_FONT
    run.font.size = Pt(10)
    run.font.color.rgb = LIGHT_GRAY


def add_body_box(slide, left=Inches(0.5), top=Inches(1.15),
                 width=None, height=None):
    if width is None:
        width = prs.slide_width - Inches(1.0)
    if height is None:
        height = prs.slide_height - Inches(1.65)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def style_run(run, size=18, bold=False, color=DARK_GRAY, mono=False,
              italic=False):
    run.font.name = MONO_FONT if mono else BODY_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_paragraph(tf, text, *, size=18, bold=False, color=DARK_GRAY,
                  mono=False, italic=False, bullet=False, indent=0,
                  space_after=4, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = indent
    p.space_after = Pt(space_after)
    prefix = ""
    if bullet:
        prefix = "•  " if indent == 0 else "–  "
    run = p.add_run()
    run.text = prefix + text
    style_run(run, size=size, bold=bold, color=color, mono=mono, italic=italic)
    return p


def add_rich_paragraph(tf, runs, *, indent=0, space_after=4, first=False,
                       bullet=False):
    """runs = list of (text, dict-of-style-kwargs)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = indent
    p.space_after = Pt(space_after)
    if bullet:
        bullet_run = p.add_run()
        bullet_run.text = "•  " if indent == 0 else "–  "
        style_run(bullet_run, size=runs[0][1].get("size", 18),
                  color=runs[0][1].get("color", DARK_GRAY))
    for text, kw in runs:
        r = p.add_run()
        r.text = text
        style_run(r, **{"size": 18, "color": DARK_GRAY, **kw})
    return p


def add_table(slide, rows, *, left, top, width, height,
              header=True, col_widths=None, font_size=14, header_size=15):
    n_rows = len(rows)
    n_cols = len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(width * w / total)
    for r_idx, row in enumerate(rows):
        for c_idx, cell_val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            run = tf.paragraphs[0].add_run()
            run.text = str(cell_val)
            run.font.name = BODY_FONT
            if header and r_idx == 0:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HDR
            else:
                run.font.size = Pt(font_size)
                run.font.color.rgb = DARK_GRAY
                cell.fill.solid()
                if r_idx % 2 == 0:
                    cell.fill.fore_color.rgb = TABLE_ALT
                else:
                    cell.fill.fore_color.rgb = WHITE
    return table


def add_code_block(slide, lines, *, left, top, width, height, font_size=13):
    """Monospaced text inside a light-gray rounded rect."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.line.color.rgb = LIGHT_GRAY
    box.line.width = Pt(0.5)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF6, 0xF6, 0xF8)
    tf = box.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = MONO_FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_GRAY


# --- slides --------------------------------------------------------------
slides_built = []


def new_slide(title=None):
    slide = prs.slides.add_slide(BLANK)
    # white background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    # move to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    if title is not None:
        add_title_bar(slide, title)
    slides_built.append(slide)
    return slide


# ===== Slide 1: Cover =====
s = new_slide()
# Big accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0),
                          prs.slide_width, Inches(2.8))
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY

tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2),
                           prs.slide_width - Inches(1.6), Inches(2.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "Selective Noise Cancelling"
r.font.name = BODY_FONT
r.font.size = Pt(54)
r.font.bold = True
r.font.color.rgb = WHITE

p = tf.add_paragraph()
p.alignment = PP_ALIGN.LEFT
p.space_before = Pt(10)
r = p.add_run()
r.text = "Query-Conditioned Source Separation ile Seçici Ses Temizleme"
r.font.name = BODY_FONT
r.font.size = Pt(24)
r.font.color.rgb = WHITE

# Subtitle below the bar
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.2),
                           prs.slide_width - Inches(1.6), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Bitirme Projesi — v2.x İterasyon Raporu"
r.font.name = BODY_FONT
r.font.size = Pt(22)
r.font.bold = True
r.font.color.rgb = DARK_GRAY

p = tf.add_paragraph()
p.space_before = Pt(6)
r = p.add_run()
r.text = "Mayıs 2026 sonu durumu"
r.font.name = BODY_FONT
r.font.size = Pt(18)
r.font.italic = True
r.font.color.rgb = LIGHT_GRAY


# ===== Slide 2: Tek Cümlede Hedef =====
s = new_slide("1. Tek Cümlede Hedef")
tf = add_body_box(s)
add_rich_paragraph(tf, [
    ("Kullanıcı bir ses ya da video yükler → uygulama içinde ", {"size": 22}),
    ("hangi tür seslerin", {"size": 22, "bold": True, "color": ORANGE}),
    (" geçtiğini tespit eder → kullanıcı silmek istediklerini işaretler "
     "→ dosya temizlenmiş halde geri döner.", {"size": 22}),
], first=True, space_after=18)

add_paragraph(tf, "", size=8)
add_paragraph(tf,
    "Mart raporundan farklı olarak artık ürüne dönüşmüş bir Gradio web "
    "uygulaması mevcut.", size=18, space_after=10)
add_paragraph(tf,
    "Bu sunum, web uygulamasına geçişten itibaren (v2.0) bugünkü v2.4 "
    "planına kadar yapılan tüm adımları, neyi neden yaptığımızı ve "
    "metriklerin nasıl evrildiğini anlatıyor.", size=18)


# ===== Slide 3: Mimari Kuş Bakışı =====
s = new_slide("2. Genel Mimari — Kuş Bakışı")
diagram_lines = [
    "          Ses / video dosyası",
    "                  │",
    "                  ▼",
    "         ┌─────────────────────┐",
    "         │  webapp.py (Gradio) │",
    "         └─────────┬───────────┘",
    "                   │",
    "         ┌─────────▼────────┐         Hangi sınıflar",
    "         │  DETECT          │  ───►   geçiyor?",
    "         │  (her sınıf için │",
    "         │   maskeyi skorla)│         (Kullanıcı seçer)",
    "         └─────────┬────────┘                │",
    "                   │            ◄────────────┘",
    "         ┌─────────▼────────┐",
    "         │  REMOVE          │  ───►   Temizlenmiş ses",
    "         │  (seçilen        │         + (varsa) video",
    "         │   sınıfları      │",
    "         │   maskele)       │",
    "         └──────────────────┘",
]
add_code_block(s, diagram_lines, left=Inches(0.8), top=Inches(1.1),
                width=Inches(7.5), height=Inches(5.6), font_size=12)

# Right side annotation
tb = s.shapes.add_textbox(Inches(8.6), Inches(1.3),
                           Inches(4.4), Inches(5.4))
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "Tek bir FiLM-koşullu U-Net,",
              size=17, bold=True, color=NAVY, first=True, space_after=6)
add_paragraph(tf,
    "\"bana X sınıfını ayır\" sorgusunu cevaplayan tek elemanlı bir araç "
    "olarak hem detection hem removal'a hizmet ediyor.", size=16,
    space_after=16)
add_paragraph(tf, "Aynı model:", size=17, bold=True, color=NAVY,
              space_after=6)
add_paragraph(tf, "Detection — N kez çağrılır (her sınıf için bir kez).",
              size=15, bullet=True, space_after=4)
add_paragraph(tf, "Removal — yalnızca seçilen sınıflar için çağrılır.",
              size=15, bullet=True)


# ===== Slide 4: Neden Query-Conditioned =====
s = new_slide("3. Neden \"Query-Conditioned\" Yaklaşım?")
tf = add_body_box(s)
add_paragraph(tf, "Klasik alternatif:", size=18, bold=True, color=NAVY,
              first=True, space_after=4)
add_paragraph(tf,
    "Her sınıf için ayrı bir model, veya N çıkışlı sabit-kafa bir model.",
    size=17, bullet=True, space_after=10)

add_paragraph(tf, "Sorun:", size=18, bold=True, color=RED, space_after=4)
add_paragraph(tf,
    "ESC-50 + UrbanSound8K + FSD50K = onlarca / yüzlerce sınıf. "
    "Sınıf başına model imkânsız; sabit N-çıkışlı model her yeni sınıfta "
    "yeniden eğitim gerektirir.", size=17, bullet=True, space_after=10)

add_paragraph(tf, "Bizim çözüm:", size=18, bold=True, color=GREEN,
              space_after=4)
add_paragraph(tf,
    "Tek bir U-Net, ek bir girişten (one-hot class query) hangi sınıfı "
    "çıkaracağını öğrenir. FiLM (Feature-wise Linear Modulation) "
    "katmanları, sorgu vektörünü her seviyede özellik haritalarına çevirip "
    "modülasyon parametreleri (γ, β) üretir.", size=17, bullet=True,
    space_after=10)

diagram = [
    "  log_mag  ───►  Encoder ───►  Bottleneck ───►  Decoder ───►  Mask",
    "                    ▲             ▲                ▲           (256×128)",
    "                    │             │                │",
    "       ┌────────────┴─────────────┴────────────────┘",
    "       │",
    "  class_query  ──►  Embed (128-d)  ──►  γ, β  her seviyeye",
]
add_code_block(s, diagram, left=Inches(0.7), top=Inches(5.2),
                width=Inches(12.0), height=Inches(1.7), font_size=12)


# ===== Slide 5: Spectrogram Sözleşmesi =====
s = new_slide("4. Spectrogram Sözleşmesi (Sabit Kalan Şey)")
tf = add_body_box(s, height=Inches(0.8))
add_paragraph(tf,
    "Tüm modelin gördüğü/ürettiği her şey aynı boyutlu spektrograma "
    "indirgenir. Web app, eğitim ve değerlendirme tek kontratla konuşur.",
    size=16, first=True)

add_table(s, [
    ["Parametre", "Değer", "Anlam"],
    ["SAMPLE_RATE", "16 kHz mono", "Tüm sesler bu hıza resample edilir"],
    ["N_FFT", "512", "STFT pencere boyu"],
    ["HOP_LENGTH", "128", "Pencere kayması (75% overlap)"],
    ["FREQ_BINS", "256", "Nyquist bin'i atılır, U-Net girişi"],
    ["TIME_FRAMES", "128", "~1 saniyelik chunk"],
    ["Çıkış", "(256, 128, 1)", "[0,1] aralığında soft mask"],
], left=Inches(0.7), top=Inches(2.1), width=Inches(12.0),
   height=Inches(4.0), col_widths=[2, 3, 7], font_size=15, header_size=16)

tb = s.shapes.add_textbox(Inches(0.7), Inches(6.55),
                           Inches(12.0), Inches(0.6))
tf = tb.text_frame
add_rich_paragraph(tf, [
    ("Bir U-Net çağrısı = ", {"size": 16}),
    ("~1 saniyelik ses", {"size": 16, "bold": True, "color": NAVY}),
    (". Uzun dosyalar webapp'te overlap-add ile birleştirilir.",
     {"size": 16}),
], first=True)


# ===== Slide 6: Eğitim Pipeline'ı (Mixer) =====
s = new_slide("5. Eğitim Pipeline'ı — Sentetik Karışım Üretici")
tf = add_body_box(s, height=Inches(0.7))
add_paragraph(tf,
    "Eğitim için dataset dosyası yok. SeparationMixer her adımda canlı bir "
    "örnek üretir:", size=16, first=True)

steps = [
    "1. Rastgele 1–4 sınıf seç          →  ESC-50 cat + US8K siren + ...",
    "2. Her birinden 1 sn pencere       →  rastgele amplitüd ile karıştır",
    "3. Bazen background noise ekle      →  white / pink, 15–30 dB SNR",
    "4. Bir query sınıfı seç:",
    "       • %85: mixture'da geçen bir sınıf  (pozitif örnek)",
    "       • %15: mixture'da geçmeyen sınıf    (negatif örnek → target = 0)",
    "5. ((log_mag, query, lin_mag), target_stem_mag)  döndür",
]
add_code_block(s, steps, left=Inches(0.7), top=Inches(2.0),
                width=Inches(12.0), height=Inches(2.7), font_size=14)

tb = s.shapes.add_textbox(Inches(0.7), Inches(5.0),
                           Inches(12.0), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Pozitif örnek: ", {"size": 17, "bold": True, "color": GREEN}),
    ("\"X sınıfı mixture'da var, onu çıkar\" → maske güçlü.", {"size": 17}),
], first=True, space_after=6)
add_rich_paragraph(tf, [
    ("Negatif örnek: ", {"size": 17, "bold": True, "color": RED}),
    ("\"Y sınıfı yok, çıkarmaya çalışma\" → maske sıfır.", {"size": 17}),
], space_after=10)
add_paragraph(tf,
    "Bu basit kurgu, modelin hem çıkarma hem yok-tespiti yeteneklerini "
    "aynı kayıp fonksiyonu ile öğrenmesini sağlar.", size=16)


# ===== Slide 7: Loss ve Eğitim Hedefi =====
s = new_slide("6. Loss ve Eğitim Hedefi")
tf = add_body_box(s, height=Inches(0.7))
add_paragraph(tf,
    "Çıktı bir maske olmasına rağmen, eğitim sırasında maske doğrudan "
    "mixture'a uygulanır (graph içinde Multiply) ve tahmin edilen stem "
    "magnitude üzerinden L1 alınır:", size=16, first=True)

code1 = [
    "  mask           = U-Net(log_mag, query)",
    "  estimated_stem = mask × linear_mag         ← grafiğin içinde",
    "  loss           = MultiResL1(estimated_stem, true_stem_mag)",
]
add_code_block(s, code1, left=Inches(0.7), top=Inches(2.1),
                width=Inches(12.0), height=Inches(1.4), font_size=15)

tb = s.shapes.add_textbox(Inches(0.7), Inches(3.7),
                           Inches(12.0), Inches(0.5))
tf = tb.text_frame
add_paragraph(tf, "Çok-çözünürlüklü L1 (v2.2'de eklendi):",
              size=17, bold=True, color=NAVY, first=True)

code2 = [
    "  L = L1(full) + 0.5 · L1(½-res) + 0.25 · L1(¼-res)",
]
add_code_block(s, code2, left=Inches(0.7), top=Inches(4.3),
                width=Inches(12.0), height=Inches(0.7), font_size=15)

tb = s.shapes.add_textbox(Inches(0.7), Inches(5.3),
                           Inches(12.0), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf,
    "Kaba ölçeklerde önce spektral şekli öğrenir, sonra ince detayı "
    "düzeltir → daha temiz maske, daha az \"musical noise\".",
    size=16, first=True)


# ===== Slide 8: Web App Pipeline =====
s = new_slide("7. Web Uygulaması — webapp.py")
tf = add_body_box(s, height=Inches(0.5))
add_paragraph(tf, "Detection (detect_sounds):",
              size=18, bold=True, color=NAVY, first=True)

code_det = [
    "  score = energy_ratio × (1 + CoV(mask)²)",
    "            ▲                    ▲",
    "            │                    └─ Maske ne kadar yoğunlaşmış?",
    "            └─ Maske mixture enerjisinin ne kadarını yakalıyor?",
]
add_code_block(s, code_det, left=Inches(0.7), top=Inches(1.85),
                width=Inches(12.0), height=Inches(1.5), font_size=13)

tb = s.shapes.add_textbox(Inches(0.7), Inches(3.45),
                           Inches(12.0), Inches(0.45))
tf = tb.text_frame
add_rich_paragraph(tf, [
    ("Mutlak eşik (0.05) ", {"size": 15}),
    ("ve", {"size": 15, "bold": True}),
    (" bağıl eşik (0.65 × kazanan) → \"5–15 alakasız sınıf\" sorunu kapanır.",
     {"size": 15}),
], first=True)

tb = s.shapes.add_textbox(Inches(0.7), Inches(4.05),
                           Inches(12.0), Inches(0.5))
tf = tb.text_frame
add_paragraph(tf, "Removal (remove_sounds):",
              size=18, bold=True, color=NAVY, first=True)

code_rm = [
    "  Tüm dosya için TEK STFT → faz globalde tutarlı",
    "  Her chunk (75% overlap):",
    "      amplitude_ratio = clip(est_mag / mix_mag, 0, 1)",
    "      mask           = clip(1 − strength · amplitude_ratio, 0, 1)",
    "      processed      = mix_mag × mask",
    "  Çoklu sınıf: combined *= mask_k   → Hanning OLA → ISTFT (orijinal faz)",
]
add_code_block(s, code_rm, left=Inches(0.7), top=Inches(4.65),
                width=Inches(12.0), height=Inches(2.4), font_size=13)


# ===== Slide 9: Değerlendirme =====
s = new_slide("8. Değerlendirme — Üç Otomatik Test")
add_table(s, [
    ["Script", "Ne Ölçer", "Sağlıklı Aralık"],
    ["diagnose_model.py",
     "Model çökmüş mü? FiLM sorguyu kullanıyor mu?",
     "3 / 3 PASS"],
    ["evaluate_conditioned_separator.py",
     "SI-SDRi — separation kalitesi (dB)",
     "≥ +5 dB"],
    ["evaluate_detection.py",
     "Per-class Precision / Recall / F1",
     "mean F1 ≥ 0.5"],
], left=Inches(0.7), top=Inches(1.3), width=Inches(12.0),
   height=Inches(2.6), col_widths=[4, 7, 3], font_size=15, header_size=16)

tb = s.shapes.add_textbox(Inches(0.7), Inches(4.2),
                           Inches(12.0), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Önemli: ", {"size": 17, "bold": True, "color": ORANGE}),
    ("Her checkpoint, yan dosyası ", {"size": 17}),
    ("*_classes.json", {"size": 17, "mono": True}),
    (" ile birlikte saklanıyor — query vektörünün boyutu ",
     {"size": 17}),
    ("modelin gördüğü vocab",
     {"size": 17, "bold": True}),
    (" ile sabitleniyor (FSD50K eklenince fark eden bu).", {"size": 17}),
], first=True, space_after=12)

add_paragraph(tf,
    "Bu üç skor, her eğitim sonrası tek bir Colab notebook "
    "(colab_evaluate.ipynb) ile dakikalar içinde alınıyor.", size=17)


# ===== Slide 10: Iterasyon Yolculugu (Trajectory Table) =====
s = new_slide("9. İterasyon Yolculuğu — Üst Düzey Tablo")
add_table(s, [
    ["Versiyon", "Ana Müdahale", "Sınıf", "SI-SDRi", "F1", "Durum"],
    ["v2.0", "İlk web app + agresif augmentation",
     "56", "—", "—", "Çökmüş"],
    ["v2.1", "Augmentation geri çekildi, normalize fix",
     "56", "−22.18 dB", "0.21", "Çalışıyor, pulsing var"],
    ["v2.2", "Full-encoder FiLM + multi-res loss + 75% OLA",
     "56*", "−22.79 dB", "0.13", "Pulsing gitti, yumuşak"],
    ["v2.3", "negative_prob 0.30→0.15 + FSD50K yüklendi",
     "235", "−21.49 dB", "0.02", "FSD50K FP felaketi"],
    ["v2.4", "min_clips_per_class=40 ile FSD50K kuyruk budama",
     "TBD", "pending", "pending", "Eğitim aşamasında"],
], left=Inches(0.5), top=Inches(1.2), width=Inches(12.3),
   height=Inches(4.5), col_widths=[1.3, 5.5, 1.3, 1.7, 1.1, 2.5],
   font_size=14, header_size=15)

tb = s.shapes.add_textbox(Inches(0.5), Inches(6.0),
                           Inches(12.3), Inches(1.1))
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf,
    "* v2.2'de FSD50K loader bug'ı vardı → 0 klip yüklendi, vocab 56'da kaldı.",
    size=14, italic=True, color=LIGHT_GRAY, first=True, space_after=4)
add_paragraph(tf,
    "Sonraki slaytlar her satırın hikayesini ayrı ayrı anlatıyor.",
    size=15)


# ===== Slide 11: v2.0 =====
s = new_slide("10. v2.0 — İlk Web App, Ama Çökmüş Eğitim")
tf = add_body_box(s, height=Inches(0.7))
add_paragraph(tf,
    "Bağlam: Web app ilk kez ayağa kalktı. Eğitim hiperparametreleri ise "
    "agresif tutulmuştu:",
    size=16, first=True)

add_table(s, [
    ["Parametre", "Değer", "Sonuç"],
    ["negative_prob", "0.45",
     "Örneklerin %45'i \"yok-sınıf\""],
    ["bg_noise_prob", "0.50",
     "Örneklerin %50'si gürültü içerir"],
    ["bg_snr_db_range", "(5, 20) dB",
     "Gürültü amplitüdü target'ın %56'sı"],
], left=Inches(0.5), top=Inches(2.0), width=Inches(12.3),
   height=Inches(1.8), col_widths=[3, 2.5, 6.8], font_size=15, header_size=16)

tb = s.shapes.add_textbox(Inches(0.5), Inches(4.0),
                           Inches(12.3), Inches(3.2))
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "Sonuç — üç sorun bir arada:",
              size=18, bold=True, color=RED, first=True, space_after=8)
add_paragraph(tf,
    "%45 negatif → L1 \"her zaman sıfır çıkar\" stratejisini ödüllendirir. "
    "Model \"güvenli sessiz maske\" dengesine çöker.",
    size=15, bullet=True, space_after=6)
add_paragraph(tf,
    "5 dB SNR → süpervizyon sinyali gürültünün altında, ayırma "
    "öğrenilemiyor.", size=15, bullet=True, space_after=6)
add_paragraph(tf,
    "Inference normalizasyon uyumsuzluğu: eğitim peak=1 normalize ediyor, "
    "webapp ham ses besliyordu → STFT magnitüdleri 3–10× düşük → model "
    "eğitilmemiş bölgede.", size=15, bullet=True, space_after=10)
add_rich_paragraph(tf, [
    ("Belirtiler: ", {"size": 15, "bold": True}),
    ("Detection 0–2 alakasız sınıf veriyor. Removal hiçbir sınıfta "
     "duyulabilir etki yapmıyor.", {"size": 15}),
])


# ===== Slide 12: v2.1 =====
s = new_slide("11. v2.1 — Sağlıklı Hiperparametreler + Normalize Düzeltmesi")
add_table(s, [
    ["Parametre", "v2.0", "v2.1"],
    ["negative_prob", "0.45", "0.30"],
    ["bg_noise_prob", "0.50", "0.10"],
    ["bg_snr_db_range", "(5, 20) dB", "(15, 30) dB"],
    ["Inference normalize", "yok", "peak-norm önce STFT"],
], left=Inches(0.5), top=Inches(1.2), width=Inches(6.5),
   height=Inches(2.5), col_widths=[3, 2, 3], font_size=14, header_size=15)

tb = s.shapes.add_textbox(Inches(7.3), Inches(1.2),
                           Inches(5.6), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "Model artık öğreniyor:",
              size=17, bold=True, color=GREEN, first=True, space_after=6)
add_paragraph(tf, "Diagnose: 3/3 PASS, HEALTHY",
              size=15, bullet=True, space_after=4)
add_paragraph(tf, "SI-SDRi: −22.18 dB (kötü ama anlamlı başlangıç)",
              size=15, bullet=True, space_after=4)
add_paragraph(tf, "F1: 0.21 (FP : TP ≈ 5.5 : 1)",
              size=15, bullet=True, space_after=4)
add_paragraph(tf, "Removal: duyulabilir",
              size=15, bullet=True)

# Bottom: two problems
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(4.0),
                          Inches(12.3), Inches(2.8))
box.line.color.rgb = ORANGE
box.line.width = Pt(1.5)
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFF, 0xF6, 0xEC)
tf = box.text_frame
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.word_wrap = True
add_paragraph(tf, "Ama iki gerçek dünya problemi var:",
              size=18, bold=True, color=ORANGE, first=True, space_after=8)
add_paragraph(tf,
    "① Sınır artifaktı: her ~0.25 sn'de \"tık-tık\" gibi pulsing.",
    size=16, bullet=True, space_after=6)
add_paragraph(tf,
    "② Detection çok geniş ağ atıyor — 8–10 alakasız sınıf yüzeye çıkıyor.",
    size=16, bullet=True, space_after=8)
add_paragraph(tf,
    "Bunların ikisi de inference tarafında problemler — model değil. "
    "v2.2'nin hedefi bu.", size=15, italic=True, color=DARK_GRAY)


# ===== Slide 13: v2.2 =====
s = new_slide("12. v2.2 — Mimari + Inference Cilası")
tf = add_body_box(s, height=Inches(0.6))
add_paragraph(tf, "Mimari değişiklik (model):",
              size=17, bold=True, color=NAVY, first=True, space_after=4)
add_paragraph(tf,
    "v2.1'de FiLM yalnızca bottleneck'te. v2.2'de her encoder seviyesinde "
    "(e1–e4 + bottleneck) → skip-connection'lar sınıf-özel.", size=15,
    bullet=True, space_after=4)
add_paragraph(tf,
    "Multi-resolution L1 loss: full + ½ + ¼ çözünürlük.", size=15,
    bullet=True)

add_table(s, [
    ["", "v2.1", "v2.2"],
    ["OLA adımı", "TIME_FRAMES/2 (%50 overlap)",
     "TIME_FRAMES/4 (%75 overlap)"],
    ["Detection skoru", "energy × (1 + CoV)",
     "energy × (1 + CoV²)"],
    ["Detection cutoff", "0.40 × winner", "0.65 × winner"],
], left=Inches(0.5), top=Inches(3.0), width=Inches(12.3),
   height=Inches(2.0), col_widths=[3, 4.5, 4.8], font_size=14, header_size=15)

# Result box
tb = s.shapes.add_textbox(Inches(0.5), Inches(5.3),
                           Inches(12.3), Inches(1.9))
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "Sonuç — bir kazanım, bir tuzak:",
              size=17, bold=True, color=NAVY, first=True, space_after=6)
add_rich_paragraph(tf, [
    ("✓ Pulsing artifaktı: tamamen gitti. ", {"size": 15, "bold": True,
                                                "color": GREEN}),
    ("(75% overlap çözdü.)", {"size": 15}),
], space_after=4)
add_rich_paragraph(tf, [
    ("✗ Kaldırma çok yumuşak ", {"size": 15, "bold": True, "color": RED}),
    ("— strength=1.0'da bile hedef ses fonda duyuluyor.", {"size": 15}),
], space_after=4)
add_rich_paragraph(tf, [
    ("✗ FSD50K loader bug'ı: ", {"size": 15, "bold": True, "color": RED}),
    ("hiyerarşik etiketler (\"Bark,Dog,Animal\") reddedildi → 0 klip yüklendi.",
     {"size": 15}),
])


# ===== Slide 14: v2.3 =====
s = new_slide("13. v2.3 — FSD50K Yüklendi + negative_prob Düşürüldü")
tf = add_body_box(s, height=Inches(0.5))
add_paragraph(tf, "İki değişiklik:",
              size=17, bold=True, color=NAVY, first=True, space_after=4)
add_paragraph(tf,
    "FSD50K loader fix: raw_labels[0] (yaprak) canonical sınıf olarak "
    "kullanılıyor → vocab 56 → 235.", size=15, bullet=True, space_after=4)
add_rich_paragraph(tf, [
    ("negative_prob: 0.30 → 0.15. ", {"size": 15, "bold": True}),
    ("Beklenti: pozitif sınıf cevabı daha güçlü, \"yumuşak kaldırma\" "
     "düzelir.", {"size": 15}),
], bullet=True)

add_table(s, [
    ["Metrik", "v2.2", "v2.3", "Δ"],
    ["SI-SDRi", "−22.79 dB", "−21.49 dB", "+1.30 dB (marjinal)"],
    ["Detection F1", "0.13", "0.02", "↓↓↓"],
    ["Total FP / TP", "564 / 92", "1092 / 13", "felaket"],
], left=Inches(0.5), top=Inches(3.0), width=Inches(12.3),
   height=Inches(2.2), col_widths=[3, 3, 3, 3.3], font_size=15, header_size=16)

tb = s.shapes.add_textbox(Inches(0.5), Inches(5.5),
                           Inches(12.3), Inches(1.7))
tf = tb.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("SI-SDRi marjinal iyileşti, ama detection F1 ", {"size": 18}),
    ("felaket düştü. ", {"size": 18, "bold": True, "color": RED}),
    ("Sebep bir sonraki slaytta tanılanıyor.", {"size": 18}),
], first=True)


# ===== Slide 15: v2.3 Tanı =====
s = new_slide("14. v2.3'ün Tanısı — FSD50K Aşırı-Hevesli Sınıflar")
tf = add_body_box(s, height=Inches(0.45))
add_paragraph(tf,
    "evaluate_detection.py'a eklediğimiz \"top-FP\" tablosu sebebi anında "
    "ortaya çıkardı:", size=15, first=True)

add_table(s, [
    ["Sınıf", "FP", "Yerel ses?"],
    ["purr", "43", "hayır (FSD50K-only)"],
    ["bass_guitar", "42", "hayır"],
    ["ringtone", "40", "hayır"],
    ["telephone", "32", "hayır"],
    ["thunderstorm", "29", "evet"],
    ["boom", "27", "hayır"],
    ["animal", "26", "hayır"],
    ["bass_drum", "24", "hayır"],
], left=Inches(0.5), top=Inches(1.65), width=Inches(5.5),
   height=Inches(4.5), col_widths=[3, 1.5, 4], font_size=13, header_size=14)

tb = s.shapes.add_textbox(Inches(6.3), Inches(1.65),
                           Inches(6.5), Inches(5.4))
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf,
    "Neden ESC-50 sınıfları aynı problemi yapmıyor da bunlar yapıyor?",
    size=15, bold=True, color=NAVY, first=True, space_after=8)
add_paragraph(tf,
    "ESC-50: sınıf başına 40 temiz, izole klip.",
    size=14, bullet=True, space_after=4)
add_paragraph(tf,
    "FSD50K: AudioSet leaf etiketlerinin uzun kuyruğu — bazı sınıflar "
    "yalnızca 2–10 (gürültülü, çok-etiketli) klip.",
    size=14, bullet=True, space_after=8)
add_paragraph(tf,
    "Az veri + gürültülü etiket → dağınık, ayrımcı olmayan maske → "
    "her chunk'ta yüksek energy×(1+CoV²) skoru.",
    size=14, bullet=True, space_after=10)
add_rich_paragraph(tf, [
    ("Kök sebep: ", {"size": 14, "bold": True, "color": ORANGE}),
    ("oran (5.6:1) ESC-50 ile birebir aynı. Sorun ",
     {"size": 14}),
    ("sınıf başına veri miktarı/kalitesi.",
     {"size": 14, "bold": True}),
], space_after=8)
add_paragraph(tf,
    "Manuel testle uyumlu: fan sesinde, kafe gürültüsünde, su sesinde "
    "\"bass_guitar var\" diye sürekli işaretliyordu.",
    size=14, italic=True, color=LIGHT_GRAY)


# ===== Slide 16: v2.4 =====
s = new_slide("15. v2.4 — Tek Değişkenli Düzeltme: Clip-Floor")
tf = add_body_box(s, height=Inches(0.5))
add_paragraph(tf, "Tek satırlık değişiklik (kavramsal):",
              size=17, bold=True, color=NAVY, first=True)

code = [
    "  # load_all_datasets içinde — birleştirme SONRASI",
    "  dropped = sorted(c for c, v in merged.items()",
    "                   if len(v) < min_clips_per_class)",
    "  for cls in dropped:",
    "      del merged[cls]",
]
add_code_block(s, code, left=Inches(0.5), top=Inches(2.0),
                width=Inches(12.3), height=Inches(1.7), font_size=14)

tb = s.shapes.add_textbox(Inches(0.5), Inches(3.9),
                           Inches(12.3), Inches(3.3))
tf = tb.text_frame
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Neden 40? ", {"size": 16, "bold": True, "color": NAVY}),
    ("ESC-50'nin sınıf başına klip sayısının aynısı → ESC-50/UrbanSound8K "
     "hiç etkilenmiyor; eşik yalnızca FSD50K kuyruğunu kesiyor.",
     {"size": 16}),
], first=True, space_after=8)
add_rich_paragraph(tf, [
    ("Neden post-merge? ", {"size": 16, "bold": True, "color": NAVY}),
    ("Cross-dataset alias'lar (FSD50K bark → ESC-50 dog) önce havuzlanıyor "
     "→ FSD50K'nın \"iyi\" sınıfları korunuyor.", {"size": 16}),
], space_after=8)
add_rich_paragraph(tf, [
    ("Atıf netliği: ", {"size": 16, "bold": True, "color": NAVY}),
    ("Tek değişken değişti → eğer FP düşerse veya SI-SDRi değişirse, "
     "sebebi bilebiliriz. v2.2'de iki şeyi (mimari + FSD50K bug) aynı anda "
     "değiştirdiğimiz için neyin ne kadar etki ettiğini ayıramadığımız "
     "dersini aldık.", {"size": 16}),
], space_after=8)
add_rich_paragraph(tf, [
    ("Durum: ", {"size": 16, "bold": True, "color": ORANGE}),
    ("Kod branch'te push'lu (commit b6ff06a). Colab eğitimi sırada.",
     {"size": 16}),
])


# ===== Slide 17: Metodolojik dersler =====
s = new_slide("16. Aldığımız Metodolojik Dersler")
tf = add_body_box(s)

add_paragraph(tf, "Ders 1 — Çok-değişkenli değişiklikler, atıf kaybı:",
              size=17, bold=True, color=NAVY, first=True, space_after=4)
add_paragraph(tf,
    "v2.2'de mimari (full-FiLM) + loss (multi-res) + FSD50K loader aynı "
    "anda değişti. FSD50K bug yüzünden net etki ölçülemedi.",
    size=14, bullet=True, space_after=12)

add_paragraph(tf, "Ders 2 — \"İyi\" hiperparametre bile yanlış dozda zarar verir:",
              size=17, bold=True, color=NAVY, space_after=4)
add_paragraph(tf,
    "negative_prob 0.45 (v2.0) çökerttiği için 0.30'a indi. 0.30 ise "
    "\"yumuşak kaldırma\" yaratıyordu. 0.15 düzeltti. Negatif örnekler "
    "ileri-geri salınımına en hassas knob.",
    size=14, bullet=True, space_after=12)

add_paragraph(tf, "Ders 3 — Yeni dataset = yeni risk yüzeyi:",
              size=17, bold=True, color=NAVY, space_after=4)
add_paragraph(tf,
    "FSD50K eklendiğinde \"daha çok veri = daha iyi model\" varsayımı "
    "uzun kuyruk nedeniyle ters döndü. Vocab şişerken her sınıfın "
    "yeterince beslendiğini kontrol etmediğimiz için detection patladı.",
    size=14, bullet=True, space_after=12)

add_paragraph(tf, "Ders 4 — Otomatik testler manuel gözlemi doğrulamalı:",
              size=17, bold=True, color=NAVY, space_after=4)
add_paragraph(tf,
    "\"Bass_guitar her seste çıkıyor\" manuel gözlemini "
    "evaluate_detection.py'a eklediğimiz top-FP tablosu ile "
    "sayısallaştırdık.", size=14, bullet=True)


# ===== Slide 18: Acik sorunlar =====
s = new_slide("17. Bugünkü Açık Sorunlar")
tf = add_body_box(s)

add_paragraph(tf, "1. SI-SDRi hâlâ negatif (≈ −21 dB):",
              size=18, bold=True, color=RED, first=True, space_after=4)
add_paragraph(tf,
    "U-Net + mixture-fazını yeniden kullanma (spectrogram-only) yapısal "
    "bir tavan koyuyor. Olası çözümler: complex spectrogram (faz tahmini), "
    "conv-tasnet tarzı time-domain model, hard-negative mining.",
    size=15, bullet=True, space_after=16)

add_paragraph(tf, "2. Bazı sınıflar diagnose'da zayıf (cat, chirping_birds):",
              size=18, bold=True, color=RED, space_after=4)
add_paragraph(tf,
    "v2.1'den beri −1.0× / −0.6× advantage. ESC-50 cat klipleri çok kısa "
    "ve transient — modeli zorluyor.",
    size=15, bullet=True, space_after=16)

add_paragraph(tf, "3. v2.4'ün gerçek dünya kalitesi henüz ölçülmedi:",
              size=18, bold=True, color=ORANGE, space_after=4)
add_paragraph(tf,
    "Sentetik metrikler iyileşse bile, gerçek kafe/sokak kayıtlarında "
    "removal kalitatif kalitesi ayrı bir test gerektirir.",
    size=15, bullet=True)


# ===== Slide 19: Yol haritasi =====
s = new_slide("18. Sıradaki Adımlar — Kısa Vadeli Yol Haritası")
tf = add_body_box(s)

add_paragraph(tf, "v2.4 değerlendirmesi (bu hafta):",
              size=17, bold=True, color=NAVY, first=True, space_after=4)
add_paragraph(tf, "Colab T4'te eğitim (~45–90 dk).",
              size=14, bullet=True, space_after=3)
add_paragraph(tf,
    "colab_evaluate.ipynb ile diagnose + SI-SDRi + F1 + top-FP tablosu.",
    size=14, bullet=True, space_after=3)
add_paragraph(tf,
    "En az 3 gerçek dünya örneği (kafe, trafik, hayvan sesli) — kalitatif "
    "dinleme.", size=14, bullet=True, space_after=14)

add_paragraph(tf, "Eğer FP hâlâ sürerse → v2.5 (hard cross-dataset negatifler):",
              size=17, bold=True, color=NAVY, space_after=4)
add_paragraph(tf,
    "Eğitim örneklerinin bir bölümünde query bir FSD50K sınıfı, mixture "
    "yalnızca ESC-50/UrbanSound8K içeriyor (target = 0) — model FSD50K "
    "sınıflarının \"diğer datasetlere ait\" seslerde tetiklenmemesi "
    "gerektiğini doğrudan öğrenir.",
    size=14, bullet=True, space_after=14)

add_paragraph(tf, "Eğer SI-SDRi düzelmiyorsa → v3.x yönü:",
              size=17, bold=True, color=NAVY, space_after=4)
add_paragraph(tf,
    "Complex spectrogram maskeleme (magnitude + phase ayrı maske).",
    size=14, bullet=True, space_after=3)
add_paragraph(tf,
    "SI-SDR diferansiyel loss waveform üzerinde.",
    size=14, bullet=True, space_after=3)
add_paragraph(tf,
    "Time-domain mimari (Conv-TasNet) ile karşılaştırma — tez baseline'ı.",
    size=14, bullet=True)


# ===== Slide 20: Tez konumlandirmasi =====
s = new_slide("19. Nihai Hedef ve Tez Konumlandırması")
tf = add_body_box(s, height=Inches(0.6))
add_paragraph(tf, "Ürün hedefi:",
              size=18, bold=True, color=NAVY, first=True, space_after=4)
add_paragraph(tf,
    "Tek tıklama ile çalışan, gerçek dünya ses/video dosyalarında "
    "seçici ses temizleme yapabilen bir uygulama. Klasik gürültü "
    "bastırıcılar tüm gürültüyü bastırır; bu sistem kullanıcının silmek "
    "istediğini sorar.",
    size=15, bullet=True, space_after=14)

add_paragraph(tf, "Tez konumlandırması:",
              size=18, bold=True, color=NAVY, space_after=4)
add_rich_paragraph(tf, [
    ("Katkı 1: ", {"size": 15, "bold": True}),
    ("Query-conditioning ile scalable mimari — sınıf-başına model "
     "çoğaltma probleminin elimine edilmesi.", {"size": 15}),
], bullet=True, space_after=4)
add_rich_paragraph(tf, [
    ("Katkı 2: ", {"size": 15, "bold": True}),
    ("Detection + removal pipeline'ının tek modelle yürütülmesi — "
     "energy-ratio × CoV² puanlaması.", {"size": 15}),
], bullet=True, space_after=4)
add_rich_paragraph(tf, [
    ("Katkı 3 (metodolojik): ", {"size": 15, "bold": True}),
    ("Çok-sınıflı, çok-dataset'li audio separation eğitiminin "
     "veri-dengeleme hassasiyetinin sistematik analizi (v2.3 → v2.4 deneyi).",
     {"size": 15}),
], bullet=True, space_after=18)

# Closing message
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), Inches(5.5),
                          Inches(12.3), Inches(1.5))
box.line.color.rgb = NAVY
box.line.width = Pt(1.5)
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xEC, 0xF1, 0xF8)
tf = box.text_frame
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.word_wrap = True
add_rich_paragraph(tf, [
    ("Kapanış: ", {"size": 16, "bold": True, "color": NAVY}),
    ("Sistem bir ürün olarak çalışıyor; metrik tarafında SI-SDRi tavanına "
     "kadar henüz yol var, ama her sürüm hem ölçülebilir bir adım atıyor "
     "hem de bir sonraki sürüm için kontrollü bir hipotez bırakıyor.",
     {"size": 16}),
], first=True)


# ===== EK 1: Hyperparam matrix =====
s = new_slide("EK 1: Versiyon × Hiperparametre Matrisi")
add_table(s, [
    ["Param", "v2.0", "v2.1", "v2.2", "v2.3", "v2.4"],
    ["negative_prob", "0.45", "0.30", "0.30", "0.15", "0.15"],
    ["bg_noise_prob", "0.50", "0.10", "0.10", "0.10", "0.10"],
    ["bg_snr_db_range", "(5,20)", "(15,30)", "(15,30)", "(15,30)", "(15,30)"],
    ["FiLM seviyesi", "bottleneck", "bottleneck", "tüm encoder",
     "tüm encoder", "tüm encoder"],
    ["Loss", "L1", "L1", "Multi-res L1", "Multi-res L1", "Multi-res L1"],
    ["OLA step", "—", "T/2", "T/4", "T/4", "T/4"],
    ["Det. cutoff", "—", "0.40·w", "0.65·w", "0.65·w", "0.65·w"],
    ["Det. skor", "—", "e·(1+CoV)", "e·(1+CoV²)", "e·(1+CoV²)",
     "e·(1+CoV²)"],
    ["FSD50K", "yok", "yok", "bug, 0 klip", "235 sınıf", "≥40 klip eşiği"],
    ["Inference norm", "bug", "düzeltildi", "—", "—", "—"],
], left=Inches(0.5), top=Inches(1.1), width=Inches(12.3),
   height=Inches(5.8), col_widths=[3, 1.7, 1.7, 1.9, 1.9, 2.1],
   font_size=12, header_size=13)


# ===== EK 2: Repo ve branch =====
s = new_slide("EK 2: Repo ve Branch Durumu")
tf = add_body_box(s)
add_rich_paragraph(tf, [
    ("Branch: ", {"size": 17, "bold": True}),
    ("feature/separator-quality-overhaul", {"size": 17, "mono": True}),
], first=True, bullet=True, space_after=6)
add_rich_paragraph(tf, [
    ("Son commit: ", {"size": 17, "bold": True}),
    ("56ac331 (FSD50K Drive symlink fix)", {"size": 17, "mono": True}),
], bullet=True, space_after=10)

add_paragraph(tf, "Entry points:", size=17, bold=True, color=NAVY,
              space_after=4)
add_paragraph(tf, "python src/model_training/train_conditioned_separator.py",
              size=14, bullet=True, mono=True, indent=1, space_after=3)
add_paragraph(tf, "python src/model_training/evaluate_conditioned_separator.py",
              size=14, bullet=True, mono=True, indent=1, space_after=3)
add_paragraph(tf, "python src/model_training/evaluate_detection.py",
              size=14, bullet=True, mono=True, indent=1, space_after=3)
add_paragraph(tf, "python src/model_training/diagnose_model.py",
              size=14, bullet=True, mono=True, indent=1, space_after=3)
add_paragraph(tf, "python src/application/webapp.py",
              size=14, bullet=True, mono=True, indent=1, space_after=12)

add_paragraph(tf, "Dokümantasyon:", size=17, bold=True, color=NAVY,
              space_after=4)
add_paragraph(tf,
    "docs/model_training_log.md — her sürüm için tam hiperparametre + "
    "sonuç tablosu", size=14, bullet=True, space_after=3)
add_paragraph(tf,
    "docs/may_progress_report/ — Mayıs raporu (web app öncesi/sonrası)",
    size=14, bullet=True, space_after=3)
add_paragraph(tf,
    "notebooks/colab_train_conditioned_separator.ipynb — eğitim notebook'u",
    size=14, bullet=True, space_after=3)
add_paragraph(tf,
    "notebooks/colab_evaluate.ipynb — değerlendirme notebook'u",
    size=14, bullet=True)


# ===== Final: Tesekkurler =====
s = new_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8),
                          prs.slide_width, Inches(1.9))
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
tb = s.shapes.add_textbox(Inches(0.5), Inches(3.0),
                           prs.slide_width - Inches(1.0), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Teşekkürler"
r.font.name = BODY_FONT
r.font.size = Pt(60)
r.font.bold = True
r.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.5), Inches(5.2),
                           prs.slide_width - Inches(1.0), Inches(1.0))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Sorular?"
r.font.name = BODY_FONT
r.font.size = Pt(28)
r.font.italic = True
r.font.color.rgb = DARK_GRAY


# --- Footers (after all slides built, so we know total) ------------------
total = len(slides_built)
for i, slide in enumerate(slides_built, start=1):
    add_footer(slide, i, total)


prs.save(OUT)
print(f"Wrote {OUT} ({OUT.stat().st_size / 1024:.1f} KB, {total} slides)")
